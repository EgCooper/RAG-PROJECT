"""Extracción de presentaciones PPT/PPTX para indexación narrativa."""

import os

from unstructured.cleaners.core import clean_extra_whitespace

from src.ingestion.cleaning import limpiar_elementos


def _limpiar_texto(texto):
    return clean_extra_whitespace(texto or "").strip()


def _elementos_desde_slides(slides):
    """Convierte slides de python-pptx a objetos compatibles con chunker."""
    from types import SimpleNamespace

    elementos = []
    for numero, slide in enumerate(slides, start=1):
        partes = []
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            texto = _limpiar_texto(shape.text)
            if texto:
                partes.append(texto)

        if not partes:
            continue

        texto_slide = "\n".join(partes)
        meta = SimpleNamespace(page_number=numero)
        elem = SimpleNamespace(text=texto_slide, metadata=meta)
        elementos.append(elem)

    return elementos


def _extraer_pptx(ruta):
    from pptx import Presentation

    prs = Presentation(ruta)
    return _elementos_desde_slides(prs.slides)


def _extraer_ppt(ruta):
    try:
        from unstructured.partition.ppt import partition_ppt

        return limpiar_elementos(partition_ppt(ruta))
    except Exception as e:
        raise ValueError(
            "Formato .ppt legacy requiere LibreOffice instalado. "
            "Convierta el archivo a .pptx o instale LibreOffice."
        ) from e


def extraer_presentacion(ruta):
    extension = os.path.splitext(ruta)[1].lower()
    if extension == ".pptx":
        return _extraer_pptx(ruta)
    if extension == ".ppt":
        return _extraer_ppt(ruta)
    raise ValueError(f"Extensión de presentación no soportada: {extension}")
